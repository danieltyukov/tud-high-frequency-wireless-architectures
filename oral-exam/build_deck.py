"""Build the EE4730 oral exam presentation (Time Domain topic) on the TU Delft template.

Usage: python3 build_deck.py
Output: EE4730_Oral_Exam_TimeDomain_Daniel_Tyukov_5714699.pptx (convert to PDF with LibreOffice)
"""

import os
import re

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
TPL = os.path.join(REPO, "TU Delft - Powerpoint templates",
                   "TU Delft - Corporate Presentation-ENG_v1.5.pptx")
ASSETS = os.path.join(HERE, "assets")
FIGS = os.path.join(REPO, "lab1", "analysis", "figs")
PICS = os.path.join(REPO, "lab1", "pics")
OUT = os.path.join(HERE, "EE4730_Oral_Exam_TimeDomain_Daniel_Tyukov_5714699.pptx")

NAVY = RGBColor(0x0C, 0x23, 0x40)
CYAN = RGBColor(0x00, 0xA6, 0xD6)
BODY = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x63, 0x63, 0x63)
LIGHT = RGBColor(0xE6, 0xF6, 0xFB)

FONT = "Arial"

# ---------------------------------------------------------------- text helpers

TOKEN = re.compile(r"(_\{[^}]*\}|\^\{[^}]*\})")


def parse_runs(s):
    """Split markup into (text, sub, sup) runs. _{x} subscript, ^{x} superscript."""
    runs = []
    for part in TOKEN.split(s):
        if not part:
            continue
        if part.startswith("_{"):
            runs.append((part[2:-1], True, False))
        elif part.startswith("^{"):
            runs.append((part[2:-1], False, True))
        else:
            runs.append((part, False, False))
    return runs


def add_para(tf, text, size=15, color=BODY, bold=False, italic=False,
             bullet=False, align=PP_ALIGN.LEFT, space_after=8, line=1.04,
             first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if bullet:
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(Emu(Inches(0.22))))
        pPr.set("indent", str(-Emu(Inches(0.22))))
        text = "•  " + text
    for txt, sub, sup in parse_runs(text):
        r = p.add_run()
        r.text = txt
        r.font.name = FONT
        r.font.size = Pt(size * 0.68) if (sub or sup) else Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        if sub:
            r.font._rPr.set("baseline", "-22000")
        elif sup:
            r.font._rPr.set("baseline", "30000")
    return p


def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def caption(slide, l, t, w, text, align=PP_ALIGN.CENTER, size=10.5):
    tf = textbox(slide, l, t, w, 0.3)
    add_para(tf, text, size=size, color=GRAY, align=align, space_after=0, first=True)


# ------------------------------------------------------------- picture helpers

def img_aspect(path):
    with Image.open(path) as im:
        return im.size[0] / im.size[1]


def pic(slide, path, l, t, w=None, h=None):
    """Place picture keeping aspect; returns (l, t, w, h) in inches."""
    a = img_aspect(path)
    if w is None:
        w = h * a
    elif h is None:
        h = w / a
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    return l, t, w, h


# --------------------------------------------------------------- table helper

def styled_table(slide, l, t, w, col_w, rows):
    """rows[0] is the header. col_w in inches, must sum to w."""
    n_r, n_c = len(rows), len(rows[0])
    gf = slide.shapes.add_table(n_r, n_c, Inches(l), Inches(t), Inches(w),
                                Inches(0.32 * n_r))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if i == 0 else (
                RGBColor(0xF2, 0xF2, 0xF2) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF))
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            add_para(tf, str(val), size=12, bold=(i == 0),
                     color=RGBColor(0xFF, 0xFF, 0xFF) if i == 0 else BODY,
                     align=PP_ALIGN.LEFT, space_after=0, first=True)
    return gf


# --------------------------------------------------------------- deck assembly

prs = Presentation(TPL)
sld_ids = prs.slides._sldIdLst
for sldId in list(sld_ids):
    prs.part.drop_rel(sldId.rId)
    sld_ids.remove(sldId)

master = prs.slide_masters[0]
L_TITLE, L_PLAIN, L_CLOSE = master.slide_layouts[1], master.slide_layouts[20], master.slide_layouts[22]


def new_slide(layout, title=None):
    s = prs.slides.add_slide(layout)
    if title is not None:
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text_frame.text = title
                for p in ph.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = FONT
    return s


# ---- 1 title -----------------------------------------------------------------
s = new_slide(L_TITLE)
for ph in s.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 10:
        ph.text_frame.text = "Time-Domain Systems and THz Material Characterization"
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT
    elif idx == 11:
        ph.text_frame.text = ("EE4730 Oral Exam  |  Time-Domain Topic  |  Daniel Tyukov  |  July 2026")
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT

# ---- 2 why time domain --------------------------------------------------------
s = new_slide(L_PLAIN, "Why measure in the time domain")
tf = textbox(s, 0.79, 1.7, 6.75, 5.1)
add_para(tf, "The THz band (0.3 to 3 THz) sits above transistor f_{max} "
             "(records around 0.75 THz) and below optical sources. Generating and "
             "detecting power is the hard part.", bullet=True, first=True,
         size=16, space_after=14)
add_para(tf, "Two generation paradigms from the lectures: multiply a microwave "
             "oscillator up in frequency, or switch a DC bias with a femtosecond "
             "laser in a photoconductor.", bullet=True, size=16, space_after=14)
add_para(tf, "The photoconductive route gives a single picosecond pulse whose "
             "spectrum covers the entire band at once.", bullet=True, size=16,
         space_after=14)
add_para(tf, "Measure the field E(t), take an FFT, and one scan characterizes a "
             "material from 0.15 to 3 THz. That is what this lab does.",
         bullet=True, size=16, space_after=14)
styled_table(s, 7.95, 1.7, 4.4, [1.30, 1.55, 1.55], [
    ["", "Heterodyne", "Photoconductive"],
    ["Source", "LO chain, xN", "fs laser + DC bias"],
    ["Signal", "CW tone", "ps pulse"],
    ["Bandwidth", "~30 GHz", "up to ~2 THz"],
    ["Use here", "VNA labs", "this lab"],
])
caption(s, 7.95, 3.6, 4.4, "Lecture 2: the two ways to make THz signals")
_, _, w, h = pic(s, os.path.join(ASSETS, "tds_schematic.png"), 8.55, 4.35, w=3.2)
caption(s, 8.55, 4.35 + h + 0.08, w, "TDS spectrometer chain (lab manual)")

# ---- 3 PCA --------------------------------------------------------------------
s = new_slide(L_PLAIN, "The photoconductive antenna")
tf = textbox(s, 0.79, 1.5, 6.4, 5.3)
add_para(tf, "An LT-GaAs gap between two electrodes. The 1.59 eV band gap matches "
             "the 780 nm pump laser.", bullet=True, first=True)
add_para(tf, "The femtosecond pulse creates carriers, the DC bias accelerates "
             "them, and recombination cuts the current off after about 1 ps. "
             "That current burst radiates the THz pulse.", bullet=True)
add_para(tf, "The lecture photocurrent model has exactly three ingredients: "
             "carrier generation (Gaussian laser pulse), Drude acceleration, "
             "and exponential recombination with τ_{c}.", bullet=True)
add_para(tf, "The receiver is the same device without bias. The incident THz "
             "field accelerates the carriers, so the gated current is "
             "proportional to the field at the moment of illumination.", bullet=True)
add_para(tf, "At high optical power the radiated field screens the bias and the "
             "photocurrent saturates. That is why high-power sources use arrays "
             "of matched PCAs.", bullet=True)
_, _, w, h = pic(s, os.path.join(ASSETS, "pca_modelling.png"), 7.55, 1.6, w=4.85)
caption(s, 7.55, 1.6 + h + 0.08, w, "PCA model and Norton equivalent (Lecture 2)")

# ---- 4 stroboscopic sampling --------------------------------------------------
s = new_slide(L_PLAIN, "Stroboscopic sampling")
tf = textbox(s, 0.79, 1.42, 11.76, 1.95)
add_para(tf, "No ADC samples a THz waveform directly; electronics stop at a few "
             "gigasamples per second.", bullet=True, size=14.5, first=True)
add_para(tf, "The laser fires an identical pulse every 10 ns (100 MHz repetition). "
             "The receiver gap conducts only during the picosecond it is "
             "illuminated, so each repetition samples the field at one delay.",
         bullet=True, size=14.5)
add_para(tf, "A motorized mirror sweeps the delay (1 mm of travel adds 6.67 ps of "
             "round trip). The readout measures a DC average current "
             "i_{rec}(Δt) proportional to the field at the gate instant; "
             "sweeping Δt reconstructs E(t).", bullet=True, size=14.5)
_, _, w, h = pic(s, os.path.join(ASSETS, "tds_system.png"), 1.85, 3.42, w=9.6)
caption(s, 1.85, 3.42 + h + 0.06, w, "Full TDS chain: one laser, a beam splitter, and a delay line (Lecture 2)")

# ---- 5 SNR budget --------------------------------------------------------------
s = new_slide(L_PLAIN, "Noise: what the lecture predicts")
tf = textbox(s, 0.79, 1.75, 7.3, 5.1)
add_para(tf, "Average transmit power 1 mW at 100 MHz repetition: 10 pJ per pulse, "
             "roughly 20 W of peak power.", bullet=True, first=True, size=15.5,
         space_after=13)
add_para(tf, "With link efficiency near 1/10, the receiver current burst is about "
             "0.14 A for half a picosecond.", bullet=True, size=15.5, space_after=13)
add_para(tf, "The duty cycle dilutes it: i_{ave} = i_{peak} · τ_{rec}/T_{rep} "
             "≈ 7 µA of DC current at the amplifier input.", bullet=True,
         size=15.5, space_after=13)
add_para(tf, "The readout amplifier noise is about 70 pA, so a single sweep "
             "already delivers a signal-to-noise ratio near 100 dB.", bullet=True,
         size=15.5, space_after=13)
add_para(tf, "Averaging N sweeps beats down white noise only:", bullet=True,
         size=15.5, space_after=13)
add_para(tf, "SNR_{N} = SNR_{1} + 10·log_{10}N", size=18, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER, space_after=0)
box = s.shapes.add_shape(5, Inches(8.6), Inches(2.3), Inches(3.5), Inches(2.9))  # 5 = rounded rect
box.fill.solid()
box.fill.fore_color.rgb = LIGHT
box.line.color.rgb = CYAN
box.line.width = Pt(1.2)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.22)
tf.margin_right = Inches(0.22)
tf.margin_top = Inches(0.16)
add_para(tf, "Prediction for the lab", size=14, bold=True, color=NAVY,
         space_after=10, first=True)
add_para(tf, "N = 50   →  +17.0 dB", size=14, color=BODY, space_after=6)
add_para(tf, "N = 100  →  +20.0 dB", size=14, color=BODY, space_after=6)
add_para(tf, "N = 1000 →  +30.0 dB", size=14, color=BODY, space_after=10)
add_para(tf, "Assignment 1 tests exactly this.", size=13, italic=True, color=GRAY,
         space_after=0)

# ---- 6 lab bench ----------------------------------------------------------------
s = new_slide(L_PLAIN, "The lab bench")
p1 = os.path.join(ASSETS, "photo_beam_path_strip.jpg")
p2 = os.path.join(ASSETS, "photo_sample_pinhole_crop.jpg")
p3 = os.path.join(ASSETS, "photo_emitter_crop.jpg")
_, _, w1, h1 = pic(s, p1, 0.79, 1.45, h=3.35)
x2 = 0.79 + w1 + 0.2
_, _, w2, h2 = pic(s, p2, x2, 1.45, h=h1)
x3 = x2 + w2 + 0.2
_, _, w3, h3 = pic(s, p3, x3, 1.45, h=h1)
caption(s, 0.79, 1.45 + h1 + 0.06, w1, "Emitter (right), four TPX50 lenses, detector and LCA-S10 amplifier (left)")
caption(s, x2, 1.45 + h2 + 0.06, w2, "Slab at the pinhole focus")
caption(s, x3, 1.45 + h3 + 0.06, w3, "Emitter PCA on its carrier")
tf = textbox(s, 0.79, 5.42, 11.76, 1.6)
add_para(tf, "Both PCAs are fibre-coupled to the Menlo unit. Lenses 1 and 4 "
             "collimate at the antennas, lenses 2 and 3 form an intermediate "
             "focus where the samples go (carriers at 35 / 70 / 190 / 245 / 305 / "
             "400 / 435 mm).", bullet=True, size=13.5, space_after=5, first=True)
add_para(tf, "Alignment by maximizing detector current lens by lens; the aligned "
             "pulse peaks at 2.89 a.u.", bullet=True, size=13.5, space_after=5)
add_para(tf, "Every scan: 100 ps delay window, 33.3 fs steps, Tukey window "
             "(α = 0.1) in ScanControl.", bullet=True, size=13.5, space_after=5)

# ---- 7 averaging results --------------------------------------------------------
s = new_slide(L_PLAIN, "Averaging, measured against the prediction")
_, _, w, h = pic(s, os.path.join(FIGS, "fig_a1_spectra.png"), 1.2, 1.5, w=10.9)
tf = textbox(s, 0.79, 1.5 + h + 0.18, 11.76, 1.6)
add_para(tf, "Dynamic range grows from 59.6 dB (N = 1) to 89.3 dB (N = 1000): "
             "gains of 16.9 / 19.3 / 29.7 dB against the predicted 17 / 20 / 30 dB. "
             "The white-noise law holds within 0.7 dB.", bullet=True, size=14.5,
         first=True)
add_para(tf, "The time trace saturates: past N = 100 the pre-pulse floor stays "
             "near 4·10^{-4} of the peak because a coherent baseline "
             "structure repeats in every sweep. Averaging only buys bandwidth "
             "where the noise is white.", bullet=True, size=14.5)

# ---- 8 extraction method --------------------------------------------------------
s = new_slide(L_PLAIN, "Extracting permittivity and loss")
tf = textbox(s, 0.79, 1.42, 11.76, 3.1)
add_para(tf, "Differential measurement: a reference scan without the sample, a "
             "sample scan with it, and their ratio H(f) = E_{sam}(f) / E_{ref}(f). "
             "The N = 100 trace is the reference.", bullet=True, size=14.5, first=True)
add_para(tf, "Delay gives the permittivity: the slab replaces a thickness d of "
             "air, so the pulse arrives Δτ = (n−1)·d/c_{0} "
             "later and ε_{r} = n^{2}.", bullet=True, size=14.5)
add_para(tf, "Loss comes from the spectrum: model the slab as a transmission-line "
             "section between air lines, then fit tan δ(f) until the modelled "
             "magnitude matches the measurement.", bullet=True, size=14.5)
add_para(tf, "The model contains the Fabry-Perot etalon of the slab faces. If the "
             "internal echo lands outside the 100 ps window (slabs 1 and 3), only "
             "the first transit is kept.", bullet=True, size=14.5)
_, _, w, h = pic(s, os.path.join(ASSETS, "tl_model.png"), 2.47, 4.62, w=8.4)
caption(s, 2.47, 4.62 + h + 0.08, w, "Air / slab / air transmission-line model from the processing notes")

# ---- 9 slab results -------------------------------------------------------------
s = new_slide(L_PLAIN, "Six slabs identified")
_, _, w, h = pic(s, os.path.join(FIGS, "fig_a2_time.png"), 0.65, 1.6, w=6.15)
caption(s, 0.65, 1.6 + h + 0.07, w, "Reference and slab pulses; the delay carries ε_{r}")
styled_table(s, 7.15, 1.6, 5.4, [0.75, 1.05, 0.85, 1.20, 1.55], [
    ["Slab", "d (mm)", "ε_{r}", "tan δ (1 THz)", "Material"],
    ["1", "7.5", "1.69", "0.0010", "porous PTFE *"],
    ["2", "3.08", "1.31", "0.0026", "Gore-Tex"],
    ["3", "4.8", "7.66", "0.0131", "Ceramic"],
    ["4", "0.525", "11.72", "0.0138", "Silicon"],
    ["5", "0.525", "11.66", "0.0028", "Silicon"],
    ["6", "2.04", "2.05", "0.0082", "Teflon"],
])
tf = textbox(s, 7.15, 4.35, 5.4, 2.2)
add_para(tf, "Five of six match the manual's reference table within 3.1%.",
         bullet=True, size=13.5, first=True)
add_para(tf, "Slab 3 arrives 28 ps late and heavily attenuated; the silicon "
             "wafers show their etalon echo 12 ps after the main pulse.",
         bullet=True, size=13.5)
add_para(tf, "* Slab 1 (ε_{r} = 1.69) sits between Gore-Tex and Teflon: "
             "consistent with porous PTFE of intermediate density.",
         bullet=True, size=13.5)

# ---- 10 H fits ------------------------------------------------------------------
s = new_slide(L_PLAIN, "Transfer function fits")
_, _, w, h = pic(s, os.path.join(FIGS, "fig_a2_Hfits.png"), 2.55, 1.42, h=4.7)
tf = textbox(s, 0.79, 6.25, 11.76, 0.6)
add_para(tf, "The fit reproduces level and etalon ripple; on silicon the ripple "
             "period c_{0}/(2nd) ≈ 83 GHz pins the optical thickness.",
         bullet=True, size=13.5, align=PP_ALIGN.CENTER, first=True)

# ---- 11 slab 2 ------------------------------------------------------------------
s = new_slide(L_PLAIN, "The slab 2 thickness conflict")
_, _, w, h = pic(s, os.path.join(FIGS, "fig_slab2_dscan.png"), 0.97, 1.45, w=11.4)
tf = textbox(s, 0.79, 1.45 + h + 0.22, 11.76, 2.0)
add_para(tf, "The data file says 0.8 mm, the label on the holder says 3.08 mm. "
             "The delay alone cannot decide: both give a plausible material "
             "(ε_{r} = 2.44, HDPE-like, versus 1.31, Gore-Tex-like).",
         bullet=True, size=14.5, first=True)
add_para(tf, "Three independent tests reject 0.8 mm: the fit residual is three "
             "times worse, the required ±10% ripple at 120 GHz is absent, "
             "and the predicted internal echo at +8.3 ps (0.13 a.u.) does not "
             "exist in the residual.", bullet=True, size=14.5)
add_para(tf, "Conclusion: d = 3.08 mm and the slab is Gore-Tex. The file name "
             "recorded a different sample or a typo.", bullet=True, size=14.5)

# ---- 12 silicon wafers ----------------------------------------------------------
s = new_slide(L_PLAIN, "Silicon wafers: same index, different loss")
_, _, w, h = pic(s, os.path.join(FIGS, "fig_a2_tand_col.png"), 0.72, 1.7, w=5.5)
caption(s, 0.72, 1.7 + h + 0.07, w, "Extracted tan δ(f): silicon falls as 1/f, polymers rise")
tf = textbox(s, 6.6, 1.6, 5.95, 5.2)
add_para(tf, "Both wafers give n ≈ 3.42, within 0.3% of the literature value "
             "for high-resistivity silicon. The etalon ripple period confirms "
             "the permittivity independently.", bullet=True, size=14.5, first=True)
add_para(tf, "Their loss differs by an order of magnitude. Conduction loss "
             "follows tan δ = σ/(ω·ε_{0}·ε_{r}), "
             "so constant conductivity appears as a 1/f slope. That is the "
             "measured behaviour below 1.5 THz.", bullet=True, size=14.5)
add_para(tf, "Converting the median conductivity of wafer 5 to resistivity gives "
             "37 Ω·cm. The wafer is labelled 33 Ω·cm from a "
             "four-point-probe measurement: an independent end-to-end check of "
             "the loss extraction.", bullet=True, size=14.5)
add_para(tf, "Wafer 4 comes out at 9 Ω·cm, the more doped and lossier "
             "twin.", bullet=True, size=14.5)

# ---- 13 takeaways ---------------------------------------------------------------
s = new_slide(L_PLAIN, "What this lab showed")
tf = textbox(s, 0.79, 1.85, 8.5, 5.0)
add_para(tf, "Stroboscopic sampling works as taught: a DC photocurrent "
             "measurement reconstructs the THz field with 33 fs steps.",
         bullet=True, size=16.5, space_after=16, first=True)
add_para(tf, "The averaging law holds: 29.7 dB measured against 30 dB predicted "
             "at N = 1000. What it cannot remove is coherent structure, and the "
             "time-domain floor shows exactly that.", bullet=True, size=16.5,
         space_after=16)
add_para(tf, "One delay reading plus a transmission-line fit identified five of "
             "six materials within 3.1% of the reference permittivities.",
         bullet=True, size=16.5, space_after=16)
add_para(tf, "The same physics settles a metadata conflict (slab 2) and "
             "reproduces a wafer's labelled resistivity within 12%.",
         bullet=True, size=16.5, space_after=16)
_, _, w, h = pic(s, os.path.join(ASSETS, "photo_slab2_crop.jpg"), 9.6, 1.85, w=3.0)
caption(s, 9.6, 1.85 + h + 0.07, w, "Slab 2, marked t = 3.08 mm")

# ---- 14 closure -----------------------------------------------------------------
s = new_slide(L_CLOSE)
for ph in s.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 10:
        ph.text_frame.text = "Thank you"
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT
    elif idx == 11:
        ph.text_frame.text = "Time-Domain topic  |  Daniel Tyukov  |  EE4730 oral exam"
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT

cp = prs.core_properties
cp.author = "Daniel Tyukov"
cp.title = "Time-Domain Systems and THz Material Characterization"
cp.subject = "EE4730 oral exam, Time-Domain topic"
cp.keywords = ""
cp.comments = ""
cp.last_modified_by = "Daniel Tyukov"

prs.save(OUT)
print("saved", OUT, round(os.path.getsize(OUT) / 1e6, 2), "MB")
