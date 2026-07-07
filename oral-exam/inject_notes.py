"""Copy speech_notes.md into the PPTX speaker-notes pane, one section per slide.

Usage: python3 inject_notes.py   (run after build_deck.py)
"""

import os
import re

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(HERE, "EE4730_Oral_Exam_TimeDomain_Daniel_Tyukov_5714699.pptx")
NOTES = os.path.join(HERE, "speech_notes.md")

text = open(NOTES).read()
sections = re.split(r"^## Slide (\d+)[^\n]*\n", text, flags=re.M)
# sections = [preamble, "1", body1, "2", body2, ...]
notes = {int(sections[i]): sections[i + 1].strip() for i in range(1, len(sections), 2)}

prs = Presentation(PPTX)
for i, slide in enumerate(prs.slides, start=1):
    body = notes.get(i)
    if not body:
        continue
    tf = slide.notes_slide.notes_text_frame
    tf.text = re.sub(r"\n(?!\n)", " ", body).replace("\n\n", "\n")

prs.save(PPTX)
print(f"notes injected into {len(notes)} slides")
